"""
빠른 문서 로더 - pdfplumber / python-docx / 텍스트 직접 읽기
Docling 대비 훨씬 빠름 (ML 모델 없이 규칙 기반 파싱)
"""
import logging
import re
from pathlib import Path

logger = logging.getLogger(__name__)

# 페이지 넘김 재결합용 — 다음 페이지가 구조(조문/별표/번호/표/불릿 등)로 시작하면 새 블록으로 둠.
_STRUCT_START = re.compile(
    r"^(제\s*\d+\s*조|부\s*칙|\[\s*별표|[○◯■□▣▷•·▪◦]|\||[-–]|\d{1,2}[.)]\s|[가-힣][.)]\s|\(\d+\)|※)"
)
# 앞 페이지가 한글로 끝나면(종결부호 없이) 단어가 이어지는 중 → 공백 없이 재결합
_HANGUL_END = re.compile(r"[가-힣]$")


class FastLoader:

    def load_text(self, file_path: str | Path) -> str:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {path}")

        suffix = path.suffix.lower()        #확장자 추출

        if suffix == ".pdf":                #PDF인지 확인
            return self._load_pdf(path)
        elif suffix == ".docx":             #DOCX인지 확인
            return self._load_docx(path)
        elif suffix in (".txt", ".md"):     #TXT / MD
            return self._load_text(path)
        elif suffix == ".pptx":             #PPTX
            return self._load_pptx(path)
        elif suffix == ".hwpx":             #HWPX (한글 XML 포맷)
            return self._load_hwpx(path)
        else:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {suffix}")

    def _load_pdf(self, path: Path) -> str:     #PDF 전용 함수
        import pdfplumber                       #PDF 텍스트 추출 라이브러리
        page_texts = []                         #페이지별 텍스트(표를 읽기순서로 끼워넣음)
        with pdfplumber.open(path) as pdf:      #PDF 열기
            for page in pdf.pages:              #페이지 반복
                page_texts.append(self._extract_page_in_order(page))

        page_texts = self._strip_running_headers(page_texts)  #반복 헤더/꼬리말 제거(재결합 전)
        result = self._join_pages(page_texts)   #페이지 경계에서 문장 이어지면 재결합
        if not result.strip():                  #빈 PDF 검사
            raise RuntimeError("PDF에서 텍스트를 추출할 수 없습니다. 스캔 문서일 수 있습니다.") #pdfplumber는 OCR이 없음
        return result

    def _extract_page_in_order(self, page) -> str:
        """한 페이지를 '읽기 순서'(세로 위치)대로 추출 — 표를 원위치에 끼워 넣는다.

        기존엔 표 밖 텍스트를 통째로 먼저 뽑고 표를 페이지 끝에 붙여서, 표가 제목에서
        떨어져 나가고(예: 출석인정 표) 페이지 넘김 문장 사이에 표가 끼어드는 문제가 있었음.
        표 '위/옆/아래'의 표 밖 텍스트를 모두 살려(pdfplumber가 표 일부만 인식해 남긴
        옆 열 텍스트 손실 방지) 표와 함께 top(세로) 순서로 이어붙여 해결."""
        try:
            tables = page.find_tables()
        except Exception:
            tables = []
        if not tables:
            return page.extract_text() or ""

        bboxes = [t.bbox for t in tables]

        def _text_between(top: float, bottom: float) -> str:
            """[top, bottom) 세로 구간의 '표 밖' 텍스트 (표 옆 텍스트 포함)."""
            if bottom - top < 1:
                return ""
            def _keep(o):
                otop = o.get("top", 0)
                if not (top <= otop < bottom):
                    return False
                x0 = o.get("x0", 0)
                return not any(bx0 <= x0 <= bx1 and bt <= otop <= bb
                               for bx0, bt, bx1, bb in bboxes)
            try:
                return page.filter(_keep).extract_text() or ""
            except Exception:
                return ""

        blocks: list[str] = []
        cur = 0.0
        page_bottom = float(page.height)
        for t in sorted(tables, key=lambda tb: tb.bbox[1]):   # top(세로) 오름차순
            btop, bbottom = float(t.bbox[1]), float(t.bbox[3])
            above = _text_between(cur, btop)                  # 표 위 텍스트
            if above.strip():
                blocks.append(above)
            beside = _text_between(btop, bbottom)             # 표와 같은 행의 표 밖 텍스트
            if beside.strip():
                # 표 옆에 텍스트가 있음 = pdfplumber가 병합셀 등으로 표를 '부분만' 인식한 것.
                # 마크다운 표로 내보내면 '행 라벨(옆 텍스트)'과 '값(표)'이 서로 다른 청크로
                # 갈라져, 값만 있고 행 이름이 없는 반쪽 청크가 생겨 검색·답변이 망가진다.
                # → 표 영역 전체를 '텍스트'로 묶어 라벨+값을 한 덩어리로 유지.
                full = self._crop_all_text(page, btop, bbottom)
                blocks.append(full if full.strip() else beside)
            else:
                md = self._safe_table_md(t)                   # 온전히 인식된 표 → 마크다운
                if md.strip():
                    blocks.append(md)
            cur = max(cur, bbottom)
        below = _text_between(cur, page_bottom)               # 마지막 표 아래 텍스트
        if below.strip():
            blocks.append(below)
        return "\n\n".join(blocks)

    def _safe_table_md(self, t) -> str:
        try:
            return self._table_to_markdown(t.extract())
        except Exception:
            return ""

    @staticmethod
    def _crop_all_text(page, top: float, bottom: float) -> str:
        """[top, bottom) 세로 밴드의 '모든' 텍스트(표 셀 포함)를 그대로 추출."""
        top = max(0.0, top - 1)
        bottom = min(float(page.height), bottom + 1)
        if bottom - top < 1:
            return ""
        try:
            return page.crop((0, top, float(page.width), bottom)).extract_text() or ""
        except Exception:
            return ""

    @staticmethod
    def _strip_running_headers(page_texts: list[str]) -> list[str]:
        """여러 페이지에 반복되는 짧은 줄(러닝 헤더/꼬리말: 규정명·문서번호 등)을 제거.
        페이지 넘김 지점에서 이런 반복 줄이 문장 사이에 끼어드는 것을 방지
        (예: '…승인으[우송대학교 규정 …[4-1-4a]]로 함)')."""
        from collections import Counter
        pages = [(pt or "") for pt in page_texts]
        nonempty = [p for p in pages if p.strip()]
        if len(nonempty) < 3:
            return page_texts        # 페이지가 적으면 반복 판단 불가 → 건드리지 않음
        cnt: Counter = Counter()
        for p in pages:
            for ln in {l.strip() for l in p.split("\n") if l.strip()}:
                cnt[ln] += 1
        thresh = max(3, int(len(nonempty) * 0.5))
        # 조문/구조 줄(제N조·부칙·번호 등)은 반복돼도 내용이므로 헤더로 오인해 지우지 않음
        boiler = {
            ln for ln, c in cnt.items()
            if c >= thresh and len(ln) <= 60 and not _STRUCT_START.match(ln)
        }
        if not boiler:
            return page_texts
        return ["\n".join(ln for ln in p.split("\n") if ln.strip() not in boiler) for p in pages]

    @staticmethod
    def _join_pages(page_texts: list[str]) -> str:
        """페이지들을 합치되, 앞 페이지가 한글 단어 중간에서(종결부호 없이) 끝나고
        다음 페이지가 구조(조문/별표/번호/표…) 시작이 아니면 → 공백 없이 이어붙여
        페이지 넘김에서 갈라진 문장을 복원한다."""
        out = ""
        for pt in page_texts:
            pt = (pt or "").strip()
            if not pt:
                continue
            if not out:
                out = pt
                continue
            prev = out.rstrip()
            first = pt.lstrip()
            if _HANGUL_END.search(prev) and not _STRUCT_START.match(first):
                out = prev + first            # 단어 연속 → 공백 없이 재결합
            else:
                out = prev + "\n\n" + pt
        return out

    @staticmethod
    def _table_to_markdown(rows: list[list]) -> str:
        """pdfplumber로 추출한 표(list[list[str|None]])를 마크다운 표로 변환.
        빈 행 제거, 셀 내 줄바꿈은 공백으로, 셀 내 파이프(|)는 대체해 표 깨짐 방지."""
        clean_rows = [r for r in (rows or []) if any((c or "").strip() for c in r)]
        if not clean_rows:
            return ""

        def fmt(c):
            return (c or "").replace("\n", " ").replace("|", "／").strip()

        ncol = max(len(r) for r in clean_rows)

        # 1열 '표'는 데이터 표가 아니라 강조 박스(테두리만 있는 안내문) →
        # 마크다운 표로 만들지 않고 평문으로 반환 (셀 전체가 article로 오염되는 것 방지)
        if ncol <= 1:
            cells = [fmt(r[0]) for r in clean_rows if r and fmt(r[0])]
            return "\n\n".join(cells)

        def pad(r):
            return list(r) + [""] * (ncol - len(r))

        # pdfplumber가 페이지 전체를 표로 잡으면 제목·섹션헤딩·불릿까지 '첫 칸만 채워진
        # 전폭(全幅) 행'으로 딸려 들어온다. 이들은 데이터가 아니라 본문이므로 평문 줄로
        # 되돌리고, 연속된 진짜 데이터 행만 마크다운 표로 만든다.
        # (안 그러면 '1. 선발 인원' 같은 헤딩이 표 행으로 청킹돼 문서가 행 단위로 파편화됨)
        # 첫 칸이 빈 연속행(| | 값 |)은 표의 이어지는 행이므로 본문으로 취급하지 않는다.
        def _is_text_row(r) -> bool:
            cells = [fmt(c) for c in pad(r)]
            filled = [i for i, c in enumerate(cells) if c]
            return len(filled) == 1 and filled[0] == 0

        out: list[str] = []
        run: list[list] = []          # 연속된 데이터 행 버퍼

        def flush():
            if not run:
                return
            if len(run) == 1:         # 데이터 행 하나뿐이면 표로 만들 것 없이 평문으로
                out.append(" ".join(c for c in (fmt(x) for x in pad(run[0])) if c))
            else:
                hdr = pad(run[0])
                out.append("| " + " | ".join(fmt(c) for c in hdr) + " |")
                out.append("| " + " | ".join("---" for _ in hdr) + " |")
                for r in run[1:]:
                    out.append("| " + " | ".join(fmt(c) for c in pad(r)) + " |")
            run.clear()

        for r in clean_rows:
            if _is_text_row(r):
                flush()
                out.append(fmt(pad(r)[0]))
            else:
                run.append(r)
        flush()
        return "\n".join(out)

    def _load_docx(self, path: Path) -> str:    #Word 전용 함수
        from docx import Document
        doc = Document(path)                    #문서 열기
        texts = [para.text for para in doc.paragraphs if para.text.strip()] #문단 추출

        # 표(Table) 안 텍스트도 추출 (문단만 읽으면 표 내용 누락됨)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text and cell_text not in texts:  # 중복 방지
                        texts.append(cell_text)

        return "\n\n".join(texts)               #합치기

    def _load_text(self, path: Path) -> str:                      #TXT / MD 전용
        return path.read_text(encoding="utf-8", errors="ignore")  #그대로 반환

    def _load_pptx(self, path: Path) -> str:    #PPT 전용
        from pptx import Presentation           #PPT 읽기 라이브러리
        prs = Presentation(path)                #PPT 열기
        texts = []
        for slide in prs.slides:                #슬라이드 반복
            for shape in slide.shapes:          #도형 반복 / PPT 안에는 텍스트박스, 제목, 표, 도형 등이 있음
                if hasattr(shape, "text") and shape.text.strip():   #텍스트 있는지 확인 / 빈 문자열 제외
                    texts.append(shape.text)    #저장
        return "\n\n".join(texts)               #합치기


    def _load_hwpx(self, path: Path) -> str:
        import zipfile
        from lxml import etree

        HP_NS = "http://www.hancom.co.kr/hwpml/2011/paragraph"
        texts = []

        with zipfile.ZipFile(path, "r") as zf:
            section_files = sorted(
                f for f in zf.namelist()
                if f.startswith("Contents/section") and f.endswith(".xml")
            )
            for section in section_files:
                root = etree.fromstring(zf.read(section))
                for t in root.iter(f"{{{HP_NS}}}t"):
                    if t.text and t.text.strip():
                        texts.append(t.text.strip())

        result = "\n".join(texts)
        if not result.strip():
            raise RuntimeError("HWPX에서 텍스트를 추출할 수 없습니다.")
        return result


    def _load_hwp(self, path: Path) -> str:
        # HWP 바이너리 전용 (pyhwp 내부 XML 변환 후 Paragraph 단위로 텍스트 추출)
        import io
        from hwp5.xmlmodel import Hwp5File
        from contextlib import closing
        import xml.etree.ElementTree as ET
        
        output = io.BytesIO()
        with closing(Hwp5File(str(path))) as hwp5file:
            hwp5file.xmlevents().dump(output)
            
        xml_data = output.getvalue()
        root = ET.fromstring(xml_data)
        
        seen = set()
        paragraphs = []
        for p in root.iter('Paragraph'):
            texts = []
            for t in p.iter('Text'):
                if t.text:
                    texts.append(t.text)
            if texts:
                line = ''.join(texts).strip()
                # HWP 레이아웃 특성상 같은 문단이 중복 저장될 수 있어 dedup 처리
                if line and line not in seen:
                    seen.add(line)
                    paragraphs.append(line)

        text = '\n'.join(paragraphs)
        
        if not text.strip():
            raise RuntimeError("HWP에서 텍스트를 추출할 수 없습니다.")
        return text


# 모듈 수준 싱글턴 — 매번 FastLoader()를 새로 만들지 않고
# `from app.rag.Loader.fast_loader import fast_loader` 후 바로 사용 가능
fast_loader = FastLoader()